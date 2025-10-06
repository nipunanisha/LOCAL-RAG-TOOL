import os
import io
import fitz  # PyMuPDF for PDF text + image extraction
from typing import Iterable, Tuple

import chromadb
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document

SUPPORTED_EXTS = {".pdf", ".txt", ".md", ".docx", ".ppt", ".pptx", ".csv", ".html"}


def load_text_from_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()

    # ---------------------- IMAGE FILES ----------------------
    if ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff"}:
        from PIL import Image
        import pytesseract

        image = Image.open(path)
        return pytesseract.image_to_string(image)

    # ---------------------- PDF FILES ----------------------
    elif ext == ".pdf":
        from PIL import Image
        import pytesseract

        text = ""
        with fitz.open(path) as doc:
            for page in doc:
                # Extract normal text
                text += (page.get_text("text") or "") + "\n"

                # Extract images and run OCR
                for img in page.get_images(full=True):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        text += "\n" + ocr_text.strip()
        return text

    # ---------------------- TEXT / MARKDOWN FILES ----------------------
    elif ext in {".txt", ".md"}:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # ---------------------- WORD DOCUMENTS ----------------------
    elif ext == ".docx":
        from PIL import Image
        import pytesseract

        doc = Document(path)
        texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        # Tables
        for table in doc.tables:
            for row in table.rows:
                texts.append(" | ".join(cell.text.strip() for cell in row.cells))

        # Embedded images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                try:
                    image = Image.open(io.BytesIO(image_data))
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        texts.append(ocr_text.strip())
                except Exception as e:
                    print(f"Error reading image from DOCX: {e}")

        return "\n".join(texts)

    # ---------------------- POWERPOINT FILES ----------------------
    elif ext in {".ppt", ".pptx"}:
        from pptx import Presentation
        from PIL import Image
        import pytesseract

        presentation = Presentation(path)
        texts = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                # Text boxes
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

                # Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        texts.append(" | ".join(cell.text.strip() for cell in row.cells))

                # Images (OCR)
                if shape.shape_type == 13:  # Picture
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        image = Image.open(image_stream)
                        ocr_text = pytesseract.image_to_string(image)
                        if ocr_text.strip():
                            texts.append(ocr_text.strip())
                    except Exception as e:
                        print(f"Error reading image from PPTX: {e}")

        return "\n".join(texts)

    # ---------------------- CSV FILES ----------------------
    elif ext == ".csv":
        import pandas as pd
        df = pd.read_csv(path, encoding="utf-8", encoding_errors="ignore")
        return df.to_csv(index=False)

    # ---------------------- HTML FILES ----------------------
    elif ext == ".html":
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
            for s in soup(["script", "style"]):
                s.extract()
            return soup.get_text(separator="\n")

    else:
        raise ValueError(f"Unsupported extension: {ext}")


# ---------------------- HELPER FUNCTIONS ----------------------
def iter_files(folder: str) -> Iterable[str]:
    for root, _, files in os.walk(folder):
        for name in files:
            if os.path.splitext(name)[1].lower() in SUPPORTED_EXTS:
                yield os.path.join(root, name)


def chunk_text(text: str, chunk_size: int = 800, chunk_overlap: int = 120) -> list[str]:
    if chunk_overlap >= chunk_size:
        raise ValueError("chunk_overlap must be smaller than chunk_size")

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - chunk_overlap
        if start < 0:
            start = 0
    return chunks


def get_embedding_model(model_name: str) -> SentenceTransformer:
    return SentenceTransformer(model_name, device="cuda")


def build_index_for_folder(folder: str, chroma_dir: str, config) -> Tuple[int, int]:
    os.makedirs(chroma_dir, exist_ok=True)
    client = chromadb.PersistentClient(path=chroma_dir)
    collection = client.get_or_create_collection(name="docs")

    emb_model = get_embedding_model(config.embedding_model)

    existing_docs = collection.get(include=["metadatas"])
    existing_ids = existing_docs.get("ids", [])
    existing_metas = existing_docs.get("metadatas", [])

    existing_paths = set()
    path_mtimes = {}
    for meta in existing_metas:
        path = meta["source"]
        existing_paths.add(path)
        path_mtimes[path] = meta.get("mtime", 0)

    current_files = list(iter_files(folder))
    current_paths = set(current_files)

    # Delete removed files
    to_delete = existing_paths - current_paths
    for path in to_delete:
        ids_to_delete = [id for id in existing_ids if id.startswith(f"{path}::")]
        if ids_to_delete:
            collection.delete(ids=ids_to_delete)

    # Add/update files
    doc_count = 0
    total_chunks = 0
    for path in current_files:
        mtime = os.path.getmtime(path)
        if path not in existing_paths or mtime > path_mtimes.get(path, 0):
            try:
                text = load_text_from_file(path)
                chunks = chunk_text(text)
                if not chunks:
                    continue
                vectors = emb_model.encode(chunks, normalize_embeddings=True).tolist()
                ids = [f"{path}::chunk::{i}" for i in range(len(chunks))]
                metas = [{"source": path, "chunk": i, "mtime": mtime} for i in range(len(chunks))]
                if path in existing_paths:
                    old_ids = [id for id in existing_ids if id.startswith(f"{path}::")]
                    collection.delete(ids=old_ids)
                collection.add(ids=ids, documents=chunks, embeddings=vectors, metadatas=metas)
                doc_count += 1
                total_chunks += len(chunks)
            except Exception as e:
                print(f"Error processing {path}: {e}")
    return doc_count, total_chunks
