"""
Per-format loaders. Each returns a list[Block] preserving structure
(page/slide/sheet/section metadata) so the chunker can be smart.
"""
from __future__ import annotations
import os
import io
import re
from typing import Iterable
from .chunking import Block

SUPPORTED_EXTS = {
    ".pdf", ".txt", ".md", ".docx", ".ppt", ".pptx",
    ".csv", ".html", ".htm", ".xlsx", ".xls",
}


# ---------------------------- PDF ----------------------------
def load_pdf(path: str, ocr_min_pixels: int = 200) -> list[Block]:
    import fitz  # PyMuPDF

    blocks: list[Block] = []
    try:
        doc = fitz.open(path)
    except Exception as e:
        print(f"[pdf] cannot open {path}: {e}")
        return []
    if doc.is_encrypted:
        if not doc.authenticate(""):
            print(f"[pdf] encrypted, skipping: {path}")
            doc.close()
            return []

    try:
        for page_idx, page in enumerate(doc):
            page_no = page_idx + 1
            section = _last_heading_on_page(page) or ""

            # 1. main text
            text = page.get_text("text") or ""
            if text.strip():
                blocks.append(Block(text=text, meta={"page": page_no, "section": section}))

            # 2. tables (PyMuPDF >=1.23)
            try:
                tables = page.find_tables()
                for t in tables:
                    md = _table_to_markdown(t.extract())
                    if md:
                        blocks.append(Block(text=md, meta={"page": page_no, "section": section, "kind": "table"}))
            except Exception:
                pass

            # 3. ocr embedded images — only if reasonably sized
            for img in page.get_images(full=True):
                try:
                    xref = img[0]
                    base = doc.extract_image(xref)
                    w, h = base.get("width", 0), base.get("height", 0)
                    if w < ocr_min_pixels or h < ocr_min_pixels:
                        continue
                    ocr = _ocr_bytes(base["image"])
                    if ocr.strip():
                        blocks.append(Block(text=ocr, meta={"page": page_no, "section": section, "kind": "ocr"}))
                except Exception as e:
                    print(f"[pdf] image ocr failed on page {page_no}: {e}")
    finally:
        doc.close()

    return blocks


def _last_heading_on_page(page) -> str:
    """Heuristic: pick the largest font-size span on the page as its heading."""
    try:
        d = page.get_text("dict")
        best = ("", 0)
        for blk in d.get("blocks", []):
            for line in blk.get("lines", []):
                for span in line.get("spans", []):
                    sz = span.get("size", 0)
                    txt = span.get("text", "").strip()
                    if sz > best[1] and 8 < len(txt) < 120:
                        best = (txt, sz)
        return best[0]
    except Exception:
        return ""


# ---------------------------- DOCX ----------------------------
def load_docx(path: str, ocr_min_pixels: int = 200) -> list[Block]:
    from docx import Document
    doc = Document(path)
    blocks: list[Block] = []

    current_section = ""
    paragraph_buffer: list[str] = []

    def flush_paras():
        if paragraph_buffer:
            text = "\n\n".join(paragraph_buffer).strip()
            if text:
                blocks.append(Block(text=text, meta={"section": current_section}))
            paragraph_buffer.clear()

    for p in doc.paragraphs:
        style = (p.style.name or "").lower()
        text = p.text.strip()
        if not text:
            continue
        if style.startswith("heading"):
            flush_paras()
            current_section = text
            blocks.append(Block(text=f"# {text}", meta={"section": current_section, "kind": "heading"}))
        else:
            paragraph_buffer.append(text)

    flush_paras()

    # tables → markdown
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        md = _table_to_markdown(rows)
        if md:
            blocks.append(Block(text=md, meta={"section": current_section, "kind": "table"}))

    # embedded images via OCR
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                blob = rel.target_part.blob
                from PIL import Image
                with Image.open(io.BytesIO(blob)) as im:
                    if im.width < ocr_min_pixels or im.height < ocr_min_pixels:
                        continue
                ocr = _ocr_bytes(blob)
                if ocr.strip():
                    blocks.append(Block(text=ocr, meta={"section": current_section, "kind": "ocr"}))
            except Exception as e:
                print(f"[docx] image ocr failed: {e}")

    return blocks


# ---------------------------- PPTX ----------------------------
def load_pptx(path: str, ocr_min_pixels: int = 200) -> list[Block]:
    from pptx import Presentation
    pres = Presentation(path)
    blocks: list[Block] = []

    for slide_idx, slide in enumerate(pres.slides):
        slide_no = slide_idx + 1
        # title
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            pass

        parts: list[str] = []
        if title:
            parts.append(f"# {title}")

        for shape in slide.shapes:
            # text frames
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t and t != title:
                    parts.append(t)
            # tables
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                md = _table_to_markdown(rows)
                if md:
                    parts.append(md)
            # images via OCR
            if getattr(shape, "shape_type", None) == 13:
                try:
                    blob = shape.image.blob
                    from PIL import Image
                    with Image.open(io.BytesIO(blob)) as im:
                        if im.width < ocr_min_pixels or im.height < ocr_min_pixels:
                            continue
                    ocr = _ocr_bytes(blob)
                    if ocr.strip():
                        parts.append(ocr)
                except Exception as e:
                    print(f"[pptx] image ocr failed on slide {slide_no}: {e}")

        # speaker notes
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"[notes] {notes}")
        except Exception:
            pass

        text = "\n\n".join(parts).strip()
        if text:
            blocks.append(Block(text=text, meta={"slide": slide_no, "section": title}))

    return blocks


# ---------------------------- XLSX / XLS ----------------------------
def load_xlsx(path: str) -> list[Block]:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    blocks: list[Block] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # detect header row — first row with all-non-None values
        header = [str(c) if c is not None else f"col_{i}" for i, c in enumerate(rows[0])]
        body = rows[1:]
        # group rows into chunks of ~30 rows so each chunk is reasonable size
        GROUP = 30
        for i in range(0, len(body), GROUP):
            group = body[i : i + GROUP]
            lines = [_format_row(header, r) for r in group if any(c is not None for c in r)]
            text = "\n".join(lines).strip()
            if text:
                blocks.append(Block(
                    text=text,
                    meta={"sheet": sheet_name, "section": sheet_name, "row_start": i + 2, "row_end": i + 1 + len(group)},
                ))
    wb.close()
    return blocks


def _format_row(header: list[str], row: tuple) -> str:
    parts = []
    for h, v in zip(header, row):
        if v is None or v == "":
            continue
        parts.append(f"{h}: {v}")
    return " · ".join(parts)


# ---------------------------- CSV ----------------------------
def load_csv(path: str) -> list[Block]:
    import pandas as pd
    df = pd.read_csv(path, encoding="utf-8", encoding_errors="ignore")
    if df.empty:
        return []
    header = [str(c) for c in df.columns]
    GROUP = 30
    blocks: list[Block] = []
    rows = df.to_dict(orient="records")
    for i in range(0, len(rows), GROUP):
        group = rows[i : i + GROUP]
        lines = [_format_row(header, tuple(r.get(h) for h in header)) for r in group]
        text = "\n".join(l for l in lines if l).strip()
        if text:
            blocks.append(Block(text=text, meta={"section": "csv", "row_start": i + 2, "row_end": i + 1 + len(group)}))
    return blocks


# ---------------------------- HTML ----------------------------
def load_html(path: str) -> list[Block]:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f, "html.parser")
    title = (soup.title.string.strip() if soup.title and soup.title.string else "")
    # nuke noise
    for tag in soup(["script", "style", "nav", "footer", "aside", "header", "noscript", "form"]):
        tag.decompose()

    blocks: list[Block] = []
    current_section = title
    buffer: list[str] = []

    def flush():
        if buffer:
            text = "\n\n".join(buffer).strip()
            if text:
                blocks.append(Block(text=text, meta={"section": current_section}))
            buffer.clear()

    body = soup.body or soup
    for el in body.descendants:
        if not getattr(el, "name", None):
            continue
        if el.name in ("h1", "h2", "h3"):
            flush()
            current_section = el.get_text(" ", strip=True)
            buffer.append(f"# {current_section}")
        elif el.name in ("p", "li", "td", "blockquote", "pre"):
            t = el.get_text(" ", strip=True)
            if t:
                buffer.append(t)
    flush()
    return blocks


# ---------------------------- TXT / MD ----------------------------
def load_text(path: str) -> list[Block]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    is_md = path.lower().endswith(".md")

    blocks: list[Block] = []
    current_section = ""
    buffer: list[str] = []

    def flush():
        if buffer:
            text = "\n\n".join(buffer).strip()
            if text:
                blocks.append(Block(text=text, meta={"section": current_section}))
            buffer.clear()

    for para in re.split(r"\n{2,}", raw):
        para = para.strip()
        if not para:
            continue
        if is_md:
            m = re.match(r"^(#{1,6})\s+(.+)$", para)
            if m:
                flush()
                current_section = m.group(2).strip()
                buffer.append(para)
                continue
        buffer.append(para)
    flush()
    return blocks


# ---------------------------- helpers ----------------------------
def _table_to_markdown(rows) -> str:
    if not rows:
        return ""
    rows = [[("" if c is None else str(c)).strip() for c in r] for r in rows]
    rows = [r for r in rows if any(r)]
    if not rows:
        return ""
    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join("---" for _ in header) + " |"]
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _ocr_bytes(data: bytes) -> str:
    from PIL import Image
    import pytesseract
    try:
        with Image.open(io.BytesIO(data)) as im:
            return pytesseract.image_to_string(im) or ""
    except Exception:
        return ""


# ---------------------------- dispatch ----------------------------
def load_blocks(path: str, ocr_min_pixels: int = 200) -> list[Block]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":              return load_pdf(path, ocr_min_pixels)
    if ext == ".docx":             return load_docx(path, ocr_min_pixels)
    if ext in (".pptx", ".ppt"):   return load_pptx(path, ocr_min_pixels)
    if ext == ".xlsx":             return load_xlsx(path)
    if ext == ".csv":              return load_csv(path)
    if ext in (".html", ".htm"):   return load_html(path)
    if ext in (".txt", ".md"):     return load_text(path)
    raise ValueError(f"Unsupported extension: {ext}")


def iter_files(folder: str) -> Iterable[str]:
    for root, _, files in os.walk(folder):
        # skip the chroma dir itself
        if os.path.basename(root) == ".chroma":
            continue
        for name in files:
            if os.path.splitext(name)[1].lower() in SUPPORTED_EXTS:
                yield os.path.join(root, name)
